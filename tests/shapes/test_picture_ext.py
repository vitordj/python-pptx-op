"""Tests for Picture.transparency added in the 'bliptransparency' branch."""

from __future__ import annotations

import io

import pytest

from pptx import Presentation
from pptx.util import Inches


def _add_picture(slide, img_path):
    return slide.shapes.add_picture(img_path, Inches(1), Inches(1), Inches(2), Inches(2))


class DescribePicture_transparency:
    """picture.transparency getter/setter via a:blip/a:alphaModFix@amt."""

    def it_returns_0_when_no_alphaModFix_is_set(self, tmp_path):
        img = _make_png(tmp_path)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = _add_picture(slide, str(img))
        assert pic.transparency == 0.0

    def it_can_set_and_read_transparency(self, tmp_path):
        img = _make_png(tmp_path)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = _add_picture(slide, str(img))
        pic.transparency = 0.5
        assert abs(pic.transparency - 0.5) < 0.001

    def it_clears_alphaModFix_when_set_to_zero(self, tmp_path):
        img = _make_png(tmp_path)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = _add_picture(slide, str(img))
        pic.transparency = 0.4
        pic.transparency = 0.0
        assert pic._blip.alphaModFix is None
        assert pic.transparency == 0.0

    def it_raises_on_out_of_range_value(self, tmp_path):
        img = _make_png(tmp_path)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = _add_picture(slide, str(img))
        with pytest.raises(ValueError, match="transparency must be between"):
            pic.transparency = 1.5

    def it_survives_a_round_trip(self, tmp_path):
        img = _make_png(tmp_path)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = _add_picture(slide, str(img))
        pic.transparency = 0.3

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)
        pic2 = prs2.slides[0].shapes[0]

        assert abs(pic2.transparency - 0.3) < 0.001


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------

def _make_png(tmp_path):
    """Create a minimal 1×1 white PNG and return its path."""
    import struct
    import zlib

    def chunk(tag, data):
        c = struct.pack(">I", len(data)) + tag + data
        return c + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)

    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
    raw = b"\x00\xff\xff\xff"
    idat = chunk(b"IDAT", zlib.compress(raw))
    iend = chunk(b"IEND", b"")

    path = tmp_path / "tiny.png"
    path.write_bytes(sig + ihdr + idat + iend)
    return path
