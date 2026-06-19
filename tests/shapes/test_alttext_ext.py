"""Tests for BaseShape.alt_text added in the 'alttext' branch."""

from __future__ import annotations

import io

import pytest

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


def _add_rect(slide):
    return slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *([Inches(1)] * 4))


class DescribeBaseShape_alt_text:
    """shape.alt_text getter/setter via p:cNvPr@descr."""

    def it_returns_empty_string_when_no_descr_is_set(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shp = _add_rect(slide)
        assert shp.alt_text == ""

    def it_can_set_and_read_alt_text(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shp = _add_rect(slide)
        shp.alt_text = "A descriptive label"
        assert shp.alt_text == "A descriptive label"

    def it_removes_the_attribute_when_assigned_empty_string(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shp = _add_rect(slide)
        shp.alt_text = "text"
        shp.alt_text = ""
        assert shp.alt_text == ""
        cNvPr = shp._element._nvXxPr.cNvPr  # pyright: ignore[reportPrivateUsage]
        assert cNvPr.descr is None

    def it_does_not_affect_the_shape_name(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shp = _add_rect(slide)
        original_name = shp.name
        shp.alt_text = "My alt text"
        assert shp.name == original_name

    def it_survives_a_round_trip(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shp = _add_rect(slide)
        shp.alt_text = "Round-trip description"

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)
        shp2 = prs2.slides[0].shapes[0]

        assert shp2.alt_text == "Round-trip description"
