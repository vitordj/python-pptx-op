"""Unit tests for LineFormat arrow heads (a:headEnd/a:tailEnd) and cap (#375)."""

import io

import pytest

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.util import Inches


def _line():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cxn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(4), Inches(1)
    )
    return cxn.line, cxn


class DescribeLineArrowHeads:
    def it_defaults_to_None(self):
        line, _ = _line()
        assert line.end_arrow.type is None
        assert line.begin_arrow.type is None

    def it_sets_tail_arrow_type_width_length(self):
        line, cxn = _line()
        line.end_arrow.type = "arrow"
        line.end_arrow.width = "lg"
        line.end_arrow.length = "lg"
        tailEnd = cxn.line._ln.find(qn("a:tailEnd"))
        assert tailEnd.get("type") == "arrow"
        assert tailEnd.get("w") == "lg"
        assert tailEnd.get("len") == "lg"
        assert line.end_arrow.type == "arrow"

    def it_sets_head_arrow_independently(self):
        line, cxn = _line()
        line.begin_arrow.type = "oval"
        assert cxn.line._ln.find(qn("a:headEnd")).get("type") == "oval"
        assert cxn.line._ln.find(qn("a:tailEnd")) is None

    def it_rejects_invalid_type(self):
        line, _ = _line()
        with pytest.raises(ValueError):
            line.end_arrow.type = "bogus"

    def it_can_remove_an_arrow_head(self):
        line, cxn = _line()
        line.end_arrow.type = "triangle"
        line.end_arrow.remove()
        assert cxn.line._ln.find(qn("a:tailEnd")) is None


class DescribeLineCap:
    def it_reads_and_writes_cap(self):
        line, cxn = _line()
        assert line.cap is None
        line.cap = "rnd"
        assert cxn.line._ln.get("cap") == "rnd"
        assert line.cap == "rnd"

    def it_rejects_invalid_cap(self):
        line, _ = _line()
        with pytest.raises(ValueError):
            line.cap = "bogus"


class DescribeRoundTrip:
    def it_survives_save_and_reopen(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cxn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(4), Inches(1)
        )
        cxn.line.end_arrow.type = "stealth"
        cxn.line.cap = "rnd"
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        cxn2 = Presentation(buf).slides[0].shapes[0]
        assert cxn2.line.end_arrow.type == "stealth"
        assert cxn2.line.cap == "rnd"
