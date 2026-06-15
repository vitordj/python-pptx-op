"""Unit tests for ColorFormat.transparency (a:alpha), issue #62."""

import pytest

from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn

from ..unitutil.cxml import element


def _color_format(cxml="a:solidFill/a:srgbClr{val=ABCDEF}"):
    return ColorFormat.from_colorchoice_parent(element(cxml))


class DescribeColorTransparency:
    def it_defaults_to_zero(self):
        cf = _color_format()
        assert cf.transparency == 0.0

    def it_can_set_transparency_writing_alpha(self):
        cf = _color_format()
        cf.transparency = 0.25
        srgbClr = cf._xFill.find(qn("a:srgbClr"))
        assert srgbClr.find(qn("a:alpha")).get("val") == "75000"
        assert cf.transparency == 0.25

    def it_removes_alpha_when_set_back_to_zero(self):
        cf = _color_format()
        cf.transparency = 0.4
        cf.transparency = 0
        assert cf._xFill.find(qn("a:srgbClr") + "/" + qn("a:alpha")) is None
        assert cf.transparency == 0.0

    def it_works_on_a_theme_color(self):
        cf = _color_format("a:solidFill/a:schemeClr{val=accent1}")
        cf.transparency = 0.5
        assert cf.type == MSO_THEME_COLOR.NOT_THEME_COLOR or cf.theme_color is not None
        assert cf._xFill.find(qn("a:schemeClr") + "/" + qn("a:alpha")).get("val") == "50000"

    @pytest.mark.parametrize("bad", [-0.1, 1.1])
    def it_rejects_out_of_range_values(self, bad):
        cf = _color_format()
        with pytest.raises(ValueError):
            cf.transparency = bad

    def it_rejects_setting_transparency_on_an_undefined_color(self):
        cf = _color_format("a:solidFill")
        with pytest.raises(ValueError):
            cf.transparency = 0.5

    def it_survives_a_round_trip_on_a_real_shape(self):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches
        import io

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *([Inches(1)] * 4))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(0x12, 0x34, 0x56)
        shp.fill.fore_color.transparency = 0.3
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)
        shp2 = prs2.slides[0].shapes[0]
        assert round(shp2.fill.fore_color.transparency, 5) == 0.3
