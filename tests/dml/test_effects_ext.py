"""Tests for ShadowFormat extended properties (blur_radius, distance, direction,
color, soft_edge_radius) added in the 'effects' branch."""

from __future__ import annotations

import io

import pytest

from pptx.dml.color import RGBColor
from pptx.dml.effect import ShadowFormat
from pptx.oxml.ns import qn
from pptx.util import Pt

from ..unitutil.cxml import element, xml


def _shadow(cxml="p:spPr"):
    return ShadowFormat(element(cxml))


class DescribeShadowFormat_shadow:
    """outer-shadow properties."""

    def it_returns_None_blur_when_no_explicit_shadow(self):
        assert _shadow().blur_radius is None

    def it_returns_None_distance_when_no_explicit_shadow(self):
        assert _shadow().distance is None

    def it_returns_None_direction_when_no_explicit_shadow(self):
        assert _shadow().direction is None

    def it_returns_None_color_when_no_explicit_shadow(self):
        assert _shadow().color is None

    def it_creates_outerShdw_with_default_subtree_on_first_write(self):
        shadow = _shadow()
        shadow.blur_radius = Pt(4)
        spPr = shadow._element
        # effectLst should be created
        assert spPr.effectLst is not None
        outerShdw = spPr.effectLst.outerShdw
        assert outerShdw is not None
        # default color child written by _new_outerShdw must be present
        assert outerShdw.find(qn("a:srgbClr")) is not None
        # the value we set should be reflected
        assert shadow.blur_radius == Pt(4)

    def it_can_set_and_read_blur_radius(self):
        shadow = _shadow()
        shadow.blur_radius = Pt(5)
        assert shadow.blur_radius == Pt(5)

    def it_can_set_and_read_distance(self):
        shadow = _shadow()
        shadow.distance = Pt(3)
        assert shadow.distance == Pt(3)

    def it_can_set_and_read_direction(self):
        shadow = _shadow()
        shadow.direction = 90.0
        assert abs(shadow.direction - 90.0) < 0.01

    def it_exposes_a_ColorFormat_for_the_shadow_color(self):
        shadow = _shadow("p:spPr/a:effectLst/a:outerShdw/a:srgbClr{val=FF0000}")
        cf = shadow.color
        assert cf is not None
        assert cf.rgb == RGBColor(0xFF, 0x00, 0x00)

    def it_can_remove_shadow_leaving_effectLst(self):
        shadow = _shadow("p:spPr/a:effectLst/a:outerShdw/a:srgbClr{val=000000}")
        shadow.remove_shadow()
        assert shadow._element.effectLst is not None
        assert shadow._element.effectLst.outerShdw is None
        assert shadow.blur_radius is None


class DescribeShadowFormat_softEdge:
    """soft-edge radius property."""

    def it_returns_None_when_no_softEdge_present(self):
        assert _shadow().soft_edge_radius is None

    def it_can_set_and_read_soft_edge_radius(self):
        shadow = _shadow()
        shadow.soft_edge_radius = Pt(6)
        assert shadow.soft_edge_radius == Pt(6)

    def it_creates_effectLst_and_softEdge_on_first_write(self):
        shadow = _shadow()
        shadow.soft_edge_radius = Pt(3)
        effectLst = shadow._element.effectLst
        assert effectLst is not None
        assert effectLst.softEdge is not None

    def it_can_clear_soft_edge_by_setting_None(self):
        shadow = _shadow()
        shadow.soft_edge_radius = Pt(5)
        shadow.soft_edge_radius = None
        assert shadow.soft_edge_radius is None

    def it_survives_a_round_trip_with_shadow_and_soft_edge(self):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *([Inches(1)] * 4))

        shadow = shp.shadow
        shadow.blur_radius = Pt(5)
        shadow.distance = Pt(3)
        shadow.direction = 90.0

        shp.shadow.soft_edge_radius = Pt(4)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)
        shp2 = prs2.slides[0].shapes[0]

        assert shp2.shadow.blur_radius == Pt(5)
        assert shp2.shadow.distance == Pt(3)
        assert abs(shp2.shadow.direction - 90.0) < 0.01
        assert shp2.shadow.soft_edge_radius == Pt(4)
