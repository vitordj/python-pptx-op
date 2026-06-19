"""Tests for _Paragraph.margin_left/margin_right/first_line_indent (paraindent branch)."""

from __future__ import annotations

import io

import pytest

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from ..unitutil.cxml import element


def _para(cxml="a:p"):
    from pptx.text.text import _Paragraph

    p = element(cxml)
    return _Paragraph(p, None)


class DescribeParagraph_margin_left:
    def it_returns_None_when_no_pPr(self):
        assert _para().margin_left is None

    def it_returns_None_when_pPr_has_no_marL(self):
        assert _para("a:p/a:pPr").margin_left is None

    def it_can_set_and_read_margin_left(self):
        para = _para()
        para.margin_left = Pt(36)
        assert para.margin_left == Pt(36)

    def it_can_clear_margin_left_by_assigning_None(self):
        para = _para()
        para.margin_left = Pt(18)
        para.margin_left = None
        assert para.margin_left is None
        assert para._p.pPr.get("marL") is None


class DescribeParagraph_margin_right:
    def it_returns_None_when_no_pPr(self):
        assert _para().margin_right is None

    def it_can_set_and_read_margin_right(self):
        para = _para()
        para.margin_right = Pt(12)
        assert para.margin_right == Pt(12)

    def it_can_clear_margin_right_by_assigning_None(self):
        para = _para()
        para.margin_right = Pt(12)
        para.margin_right = None
        assert para.margin_right is None


class DescribeParagraph_first_line_indent:
    def it_returns_None_when_no_pPr(self):
        assert _para().first_line_indent is None

    def it_can_set_a_positive_first_line_indent(self):
        para = _para()
        para.first_line_indent = Pt(18)
        assert para.first_line_indent == Pt(18)

    def it_can_set_a_negative_indent_for_hanging_effect(self):
        para = _para()
        para.first_line_indent = Pt(-18)
        assert para.first_line_indent == Pt(-18)

    def it_can_clear_first_line_indent_by_assigning_None(self):
        para = _para()
        para.first_line_indent = Pt(18)
        para.first_line_indent = None
        assert para.first_line_indent is None

    def it_survives_a_round_trip_with_all_three_properties(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shp = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        tf = shp.text_frame
        tf.text = "indented paragraph"
        para = tf.paragraphs[0]
        para.margin_left = Pt(36)
        para.margin_right = Pt(12)
        para.first_line_indent = Pt(-18)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)
        para2 = prs2.slides[0].shapes[0].text_frame.paragraphs[0]

        assert para2.margin_left == Pt(36)
        assert para2.margin_right == Pt(12)
        assert para2.first_line_indent == Pt(-18)
