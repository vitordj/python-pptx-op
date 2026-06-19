"""Tests for TextFrame.text_direction added in the 'verttext' branch."""

from __future__ import annotations

import io

import pytest

from pptx import Presentation
from pptx.enum.text import MSO_TEXT_DIRECTION
from pptx.util import Inches


def _add_textbox(slide):
    return slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))


class DescribeTextFrame_text_direction:
    """TextFrame.text_direction getter/setter via a:bodyPr@vert."""

    def it_returns_None_when_no_vert_attribute_is_set(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tf = _add_textbox(slide).text_frame
        assert tf.text_direction is None

    def it_can_set_VERTICAL_and_read_it_back(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tf = _add_textbox(slide).text_frame
        tf.text_direction = MSO_TEXT_DIRECTION.VERTICAL
        assert tf.text_direction == MSO_TEXT_DIRECTION.VERTICAL

    def it_can_set_VERTICAL_270_and_read_it_back(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tf = _add_textbox(slide).text_frame
        tf.text_direction = MSO_TEXT_DIRECTION.VERTICAL_270
        assert tf.text_direction == MSO_TEXT_DIRECTION.VERTICAL_270

    def it_can_set_HORIZONTAL_and_read_it_back(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tf = _add_textbox(slide).text_frame
        tf.text_direction = MSO_TEXT_DIRECTION.HORIZONTAL
        assert tf.text_direction == MSO_TEXT_DIRECTION.HORIZONTAL

    def it_clears_vert_when_set_to_None(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tf = _add_textbox(slide).text_frame
        tf.text_direction = MSO_TEXT_DIRECTION.VERTICAL
        tf.text_direction = None
        assert tf.text_direction is None
        assert tf._txBody.bodyPr.get("vert") is None

    def it_does_not_affect_text_content(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = _add_textbox(slide)
        shape.text_frame.text = "rotated"
        shape.text_frame.text_direction = MSO_TEXT_DIRECTION.VERTICAL
        assert shape.text_frame.text == "rotated"

    def it_survives_a_round_trip(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tf = _add_textbox(slide).text_frame
        tf.text_direction = MSO_TEXT_DIRECTION.VERTICAL

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        prs2 = Presentation(buf)
        tf2 = prs2.slides[0].shapes[0].text_frame

        assert tf2.text_direction == MSO_TEXT_DIRECTION.VERTICAL
