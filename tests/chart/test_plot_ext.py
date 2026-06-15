"""Unit tests for DoughnutPlot.hole_size / .first_slice_angle (issue #493)."""

import pytest

from pptx.chart.data import CategoryChartData
from pptx.chart.plot import DoughnutPlot
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches

from ..unitutil.cxml import element


def _doughnut_plot():
    data = CategoryChartData()
    data.categories = ["A", "B", "C"]
    data.add_series("S", (1, 2, 3))
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, Inches(1), Inches(1), Inches(5), Inches(4), data
    )
    return gframe.chart.plots[0]


class DescribeDoughnutPlotHoleSize:
    def it_defaults_to_10_when_absent(self):
        plot = DoughnutPlot(element("c:doughnutChart"), None)
        assert plot.hole_size == 10

    def it_reads_and_writes_hole_size(self):
        plot = _doughnut_plot()
        plot.hole_size = 60
        assert plot.hole_size == 60
        assert plot._element.find(qn("c:holeSize")).get("val") == "60"

    def it_reads_and_writes_first_slice_angle(self):
        plot = _doughnut_plot()
        plot.first_slice_angle = 90
        assert plot.first_slice_angle == 90
        assert plot._element.find(qn("c:firstSliceAng")).get("val") == "90"

    @pytest.mark.parametrize("bad", [0, 91, 200])
    def it_rejects_out_of_range_hole_size(self, bad):
        plot = _doughnut_plot()
        with pytest.raises(ValueError):
            plot.hole_size = bad

    @pytest.mark.parametrize("bad", [-1, 361])
    def it_rejects_out_of_range_angle(self, bad):
        plot = _doughnut_plot()
        with pytest.raises(ValueError):
            plot.first_slice_angle = bad
